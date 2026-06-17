from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor


ACCENT = RGBColor(0x1F, 0x29, 0x37)
INK = RGBColor(0x11, 0x18, 0x27)
MUTED = RGBColor(0x6B, 0x72, 0x80)


def _set_text(tf, text: str, *, size: int, bold: bool = False, color: RGBColor = INK):
    tf.text = text
    p = tf.paragraphs[0]
    for run in p.runs:
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.color.rgb = color


def build_pptx(path: str, title: str, subtitle: str, slides: list[dict]) -> None:
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    blank = prs.slide_layouts[6]

    # Title slide
    s = prs.slides.add_slide(blank)
    tx = s.shapes.add_textbox(Inches(0.8), Inches(2.6), Inches(11.7), Inches(1.6)).text_frame
    _set_text(tx, title, size=44, bold=True, color=ACCENT)
    if subtitle:
        sx = s.shapes.add_textbox(Inches(0.8), Inches(4.0), Inches(11.7), Inches(1.0)).text_frame
        _set_text(sx, subtitle, size=22, color=MUTED)

    for spec in slides:
        s = prs.slides.add_slide(blank)
        tx = s.shapes.add_textbox(Inches(0.6), Inches(0.4), Inches(12.1), Inches(0.9)).text_frame
        _set_text(tx, spec.get("title") or "", size=28, bold=True, color=ACCENT)

        bullets = spec.get("bullets") or []
        if bullets:
            bx = s.shapes.add_textbox(Inches(0.6), Inches(1.5), Inches(12.1), Inches(5.5)).text_frame
            bx.word_wrap = True
            for i, b in enumerate(bullets):
                p = bx.paragraphs[0] if i == 0 else bx.add_paragraph()
                p.text = b
                p.level = 0
                for run in p.runs:
                    run.font.size = Pt(20)
                    run.font.color.rgb = INK
                p.space_after = Pt(8)

        notes = spec.get("notes")
        if notes:
            s.notes_slide.notes_text_frame.text = notes

    prs.save(path)
